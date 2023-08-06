#!usr/bin/env python
# -*- coding: utf-8 -*-
from __future__ import unicode_literals
from __future__ import division

try:
    from xmlrpclib import ServerProxy
except ImportError:
    from xmlrpc.client import ServerProxy
import codecs
import contextlib
import logging
import datetime
import hashlib
import os
import random
import re
import json
import shlex
import shutil
import subprocess
import sys
import time
import tempfile
import traceback
import urllib
import dateparser
import pyperclip
import toml
from pptx.enum.text import PP_ALIGN

try:
    basestring
except NameError:
    basestring = str
try:
    input = raw_input
except NameError:
    pass
try:
    unquote = urllib.unquote
except AttributeError:

    def unquote(bytestring):
        return urllib.parse.unquote(bytestring.decode("utf8")).encode("utf8")


from docx import Document
from docx.shared import Inches
from PIL import Image
import pyimgur

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt

from chgksuite.common import (
    get_lastdir,
    get_chgksuite_dir,
    set_lastdir,
    on_close,
    DummyLogger,
    log_wrap,
    QUESTION_LABELS,
    check_question,
    retry_wrapper_factory,
    bring_to_front,
)
import chgksuite.typotools as typotools
from chgksuite.typotools import remove_excessive_whitespace as rew

args = None
debug = False
console_mode = False
re_url = re.compile(
    r"""((?:[a-z][\w-]+:(?:/{1,3}|[a-z0-9%])|www\d{0,3}[.]"""
    """|[a-z0-9.\-]+[.‌​][a-z]{2,4}/)(?:[^\s<>]+|(([^\s()<>]+|(([^\s<>]+)))*))+"""
    """(?:(([^\s<>]+|(‌​([^\s<>]+)))*)|[^\s`!\[\]{};:'".,<>?«»“”‘’]))""",
    re.DOTALL,
)
re_perc = re.compile(r"(%[0-9a-fA-F]{2})+")
re_scaps = re.compile(
    r"(^|[\s])([\[\]\(\)«»А-Я \u0301`ЁA-Z]{2,})([\s,!\.;:-\?]|$)"
)
re_em = re.compile(r"_(.+?)_")
re_lowercase = re.compile(r"[а-яё]")
re_uppercase = re.compile(r"[А-ЯЁ]")
re_editors = re.compile(r"^[рР]едакторы? *(пакета|тура)? *[—\-–−:] ?")

REQUIRED_LABELS = set(["question", "answer"])
IMGUR_CLIENT_ID = "8da1bd97da30ac1"
im = pyimgur.Imgur(IMGUR_CLIENT_ID)

ENC = sys.stdout.encoding or "utf8"
CONSOLE_ENC = ENC

FIELDS = {
    "zachet": "Зачёт: ",
    "nezachet": "Незачёт: ",
    "comment": "Комментарий: ",
    "source": "Источник: ",
    "author": "Автор: ",
}

WHITEN = {
    "handout": False,
    "zachet": True,
    "nezachet": True,
    "comment": True,
    "source": True,
    "author": False,
}


logger = DummyLogger()
retry_wrapper = None


def make_filename(s, ext, nots=False):
    bn = os.path.splitext(os.path.basename(s))[0]
    if nots:
        return bn + "." + ext
    return "{}_{}.{}".format(
        bn, datetime.datetime.now().strftime("%Y%m%dT%H%M"), ext
    )


@contextlib.contextmanager
def make_temp_directory(**kwargs):
    temp_dir = tempfile.mkdtemp(**kwargs)
    yield temp_dir
    shutil.rmtree(temp_dir)


def proportional_resize(tup):
    if max(tup) > 600:
        return tuple([int(x * 600 / max(tup)) for x in tup])
    if max(tup) < 200:
        return tuple([int(x * 200 / max(tup)) for x in tup])
    return tup


def imgsize(imgfile):
    img = Image.open(imgfile)
    width, height = proportional_resize((img.width, img.height))
    return width, height


def convert_size(width, height, dimensions="pixels", emsize=25, dpi=120):
    if dimensions == "pixels":
        return width, height
    if dimensions == "ems":
        return width / emsize, height / emsize
    if dimensions == "inches":
        return width / dpi, height / dpi


def search_for_imgfile(imgfile, tmp_dir, targetdir):
    if os.path.isfile(imgfile):
        return imgfile
    for dirname in [tmp_dir, targetdir]:
        if not os.path.isdir(dirname):
            continue
        imgfile2 = os.path.join(dirname, os.path.basename(imgfile))
        if os.path.isfile(imgfile2):
            return imgfile2
    raise Exception("Image file {} not found".format(imgfile))


def parse_single_size(ssize, dpi=120, emsize=25):
    if ssize.endswith("in"):
        ssize = ssize[:-2]
        return float(ssize) * dpi
    if ssize.endswith("em"):
        ssize = ssize[:-2]
        return float(ssize) * emsize
    if ssize.endswith("px"):
        ssize = ssize[:-2]
    return float(ssize)


def parseimg(s, dimensions="pixels", tmp_dir=None, targetdir=None):
    width = -1
    height = -1
    sp = s.split()
    cwd = os.getcwd()
    imgfile = sp[-1]
    imgfile = search_for_imgfile(imgfile, tmp_dir, targetdir)
    size = imgsize(imgfile)

    if len(sp) == 1:
        width, height = convert_size(*size, dimensions=dimensions)
        return imgfile.replace("\\", "/"), width, height
    else:
        for spsp in sp[:-1]:
            spspsp = spsp.split("=")
            if spspsp[0] == "w":
                width = parse_single_size(spspsp[1])
            if spspsp[0] == "h":
                height = parse_single_size(spspsp[1])
        if width != -1 and height == -1:
            height = size[1] * (width / size[0])
        elif width == -1 and height != -1:
            width = size[0] * (height / size[1])
        width, height = convert_size(width, height, dimensions=dimensions)
        return imgfile.replace("\\", "/"), width, height


def partition(alist, indices):
    return [alist[i:j] for i, j in zip([0] + indices, indices + [None])]


def parse_4s_elem(s):
    def find_next_unescaped(ss, index):
        j = index + 1
        while j < len(ss):
            if ss[j] == "\\" and j + 2 < len(ss):
                j += 2
            if ss[j] == ss[index]:
                return j
            j += 1
        return -1

    s = s.replace("\\_", "$$$$UNDERSCORE$$$$")
    for gr in re_url.finditer(s):
        gr0 = gr.group(0)
        s = s.replace(gr0, gr0.replace("_", "$$$$UNDERSCORE$$$$"))

    # for gr in re_scaps.finditer(s):
    #     gr0 = gr.group(0)
    #     s = s.replace(gr0, '(sc '+gr0.lower()+')')

    grs = sorted(
        [match.group(0) for match in re_perc.finditer(s)],
        key=len,
        reverse=True,
    )
    for gr in grs:
        try:
            s = s.replace(gr, unquote(gr.encode("utf8")).decode("utf8"))
        except:
            logger.debug(
                "error decoding on line {}: {}\n".format(
                    log_wrap(gr), traceback.format_exc()
                )
            )

    i = 0
    topart = []
    while i < len(s):
        if s[i] == "_" and (i == 0 or s[i - 1] not in {"\\", "\u6565"}):
            logger.debug("found _ at {} of line {}".format(i, s))
            topart.append(i)
            if find_next_unescaped(s, i) != -1:
                topart.append(find_next_unescaped(s, i) + 1)
                i = find_next_unescaped(s, i) + 2
                continue
        if (
            s[i] == "("
            and i + len("(img") < len(s)
            and "".join(s[i : i + len("(img")]) == "(img"
        ):
            topart.append(i)
            if typotools.find_matching_closing_bracket(s, i) is not None:
                topart.append(
                    typotools.find_matching_closing_bracket(s, i) + 1
                )
                i = typotools.find_matching_closing_bracket(s, i)
        # if (s[i] == '(' and i + len('(sc') < len(s) and ''.join(s[i:
        #                     i+len('(sc')])=='(sc'):
        #     debug_print('sc candidate')
        #     topart.append(i)
        #     if not typotools.find_matching_closing_bracket(s, i) is None:
        #         topart.append(
        #             typotools.find_matching_closing_bracket(s, i)+1)
        #         i = typotools.find_matching_closing_bracket(s, i)+2
        i += 1

    topart = sorted(topart)

    parts = [
        ["", "".join(x.replace("\u6565", ""))] for x in partition(s, topart)
    ]

    for part in parts:
        if not part[1]:
            continue
        try:
            if part[1][-1] == "_":
                part[1] = part[1][1:]
                part[0] = "em"
            if not part[1]:
                continue
            if part[1][-1] == "_":
                part[1] = part[1][:-1]
                part[0] = "em"
            if not part[1]:
                continue
            if len(part[1]) > 4 and part[1][:4] == "(img":
                if part[1][-1] != ")":
                    part[1] = part[1] + ")"
                part[1] = typotools.remove_excessive_whitespace(part[1][4:-1])
                part[0] = "img"
                logger.debug("found img at {}".format(part[1]))
            if len(part[1]) > 3 and part[1][:4] == "(sc":
                if part[1][-1] != ")":
                    part[1] = part[1] + ")"
                part[1] = typotools.remove_excessive_whitespace(part[1][3:-1])
                part[0] = "sc"
                logger.debug("found img at {}".format(log_wrap(part[1])))
            part[1] = part[1].replace("\\_", "_")
            part[1] = part[1].replace("\\.", ".")
            part[1] = part[1].replace("$$$$UNDERSCORE$$$$", "_")
        except:
            sys.stderr.write(
                "Error on part {}: {}".format(
                    log_wrap(part), traceback.format_exc()
                )
            )

    return parts


def process_list(element):
    if "-" not in element[1]:
        return
    sp = element[1].split("\n")
    sp = [rew(x) for x in sp]
    list_markers = [i for i in range(len(sp)) if sp[i].startswith("-")]
    if not list_markers or len(list_markers) == 1:
        return
    preamble = "\n".join(sp[: list_markers[0]])
    inner_list = []
    for num, index in enumerate(list_markers):
        if (num + 1) == len(list_markers):
            inner_list.append(rew("\n".join(sp[index:])[1:]))
        else:
            inner_list.append(
                rew("\n".join(sp[index : list_markers[num + 1]])[1:])
            )
    if preamble:
        element[1] = [preamble, inner_list]
    else:
        element[1] = inner_list


def parse_4s(s, randomize=False):
    mapping = {
        "#": "meta",
        "##": "section",
        "###": "heading",
        "###LJ": "ljheading",
        "#EDITOR": "editor",
        "#DATE": "date",
        "?": "question",
        "№": "number",
        "№№": "setcounter",
        "!": "answer",
        "=": "zachet",
        "!=": "nezachet",
        "^": "source",
        "/": "comment",
        "@": "author",
        ">": "handout",
    }

    structure = []

    if s[0] == "\ufeff" and len(s) > 1:
        s = s[1:]

    with codecs.open("raw.debug", "w", "utf8") as debugf:
        debugf.write(log_wrap(s.split("\n")))

    for line in s.split("\n"):
        if rew(line) == "":
            structure.append(["", ""])
        else:
            if line.split()[0] in mapping:
                structure.append(
                    [
                        mapping[line.split()[0]],
                        rew(line[len(line.split()[0]) :]),
                    ]
                )
            else:
                if len(structure) >= 1:
                    structure[len(structure) - 1][1] += "\n" + line

    final_structure = []
    current_question = {}
    counter = 1

    if debug:
        with codecs.open("debug1st.debug", "w", "utf8") as debugf:
            debugf.write(log_wrap(structure))

    for element in structure:

        # find list in element

        process_list(element)

        if element[0] in QUESTION_LABELS:
            if element[0] in current_question:

                if isinstance(
                    current_question[element[0]], basestring
                ) and isinstance(element[1], basestring):
                    current_question[element[0]] += "\n" + element[1]

                elif isinstance(
                    current_question[element[0]], list
                ) and isinstance(element[1], basestring):
                    current_question[element[0]][0] += "\n" + element[1]

                elif isinstance(
                    current_question[element[0]], basestring
                ) and isinstance(element[1], list):
                    current_question[element[0]] = [
                        element[1][0] + "\n" + current_question[element[0]],
                        element[1][1],
                    ]

                elif isinstance(
                    current_question[element[0]], list
                ) and isinstance(element[1], list):
                    current_question[element[0]][0] += "\n" + element[1][0]
                    current_question[element[0]][1] += element[1][1]
            else:
                current_question[element[0]] = element[1]

        elif element[0] == "":

            if current_question != {} and set(current_question.keys()) != {
                "setcounter"
            }:

                try:
                    assert all(
                        (True if label in current_question else False)
                        for label in REQUIRED_LABELS
                    )
                except AssertionError:
                    logger.error(
                        "Question {} misses "
                        "some of the required fields "
                        "and will therefore "
                        "be omitted.".format(log_wrap(current_question))
                    )
                    continue
                if "setcounter" in current_question:
                    counter = int(current_question["setcounter"])
                if "number" not in current_question:
                    current_question["number"] = counter
                    counter += 1
                final_structure.append(["Question", current_question])

                current_question = {}

        else:
            final_structure.append([element[0], element[1]])

    if current_question != {}:
        try:
            assert all(
                (True if label in current_question else False)
                for label in REQUIRED_LABELS
            )
            if "setcounter" in current_question:
                counter = int(current_question["setcounter"])
            if "number" not in current_question:
                current_question["number"] = counter
                counter += 1
            final_structure.append(["Question", current_question])
        except AssertionError:
            logger.error(
                "Question {} misses "
                "some of the required fields and will therefore "
                "be omitted.".format(log_wrap(current_question))
            )

    if randomize:
        random.shuffle(final_structure, lambda: 0.3)
        i = 1
        for element in final_structure:
            if element[0] == "Question":
                element[1]["number"] = i
                i += 1

    if debug:
        with codecs.open("debug.debug", "w", "utf8") as debugf:
            debugf.write(log_wrap(final_structure))

    for element in final_structure:
        if element[0] == "Question":
            check_question(element[1], logger=logger)

    return final_structure


def html_format_question(q, **kwargs):
    def yapper(x):
        return htmlyapper(x, **kwargs)

    if "setcounter" in q:
        gui_compose.counter = int(q["setcounter"])
    res = "<strong>Вопрос {}.</strong> {}".format(
        gui_compose.counter if "number" not in q else q["number"],
        yapper(q["question"])
        + ("\n<lj-spoiler>" if not args.nospoilers else ""),
    )
    if "number" not in q:
        gui_compose.counter += 1
    res += "\n<strong>Ответ: </strong>{}".format(yapper(q["answer"]))
    if "zachet" in q:
        res += "\n<strong>Зачёт: </strong>{}".format(yapper(q["zachet"]))
    if "nezachet" in q:
        res += "\n<strong>Незачёт: </strong>{}".format(yapper(q["nezachet"]))
    if "comment" in q:
        res += "\n<strong>Комментарий: </strong>{}".format(
            yapper(q["comment"])
        )
    if "source" in q:
        res += "\n<strong>Источник{}: </strong>{}".format(
            "и" if isinstance(q["source"], list) else "", yapper(q["source"])
        )
    if "author" in q:
        res += "\n<strong>Автор{}: </strong>{}".format(
            "ы" if isinstance(q["author"], list) else "", yapper(q["author"])
        )
    if not args.nospoilers:
        res += "</lj-spoiler>"
    return res


def htmlrepl(zz):
    zz = zz.replace("&", "&amp;")
    zz = zz.replace("<", "&lt;")
    zz = zz.replace(">", "&gt;")

    # while re_scaps.search(zz):
    #     zz = zz.replace(re_scaps.search(zz).group(1),
    #         '\\tsc{'+re_scaps.search(zz).group(1).lower()+'}')

    while "`" in zz:
        if zz.index("`") + 1 >= len(zz):
            zz = zz.replace("`", "")
        else:
            if zz.index("`") + 2 < len(zz) and re.search(
                r"\s", zz[zz.index("`") + 2]
            ):
                zz = zz[: zz.index("`") + 2] + "" + zz[zz.index("`") + 2 :]
            if zz.index("`") + 1 < len(zz) and re_lowercase.search(
                zz[zz.index("`") + 1]
            ):
                zz = (
                    zz[: zz.index("`") + 1]
                    + ""
                    + zz[zz.index("`") + 1]
                    + "&#x0301;"
                    + zz[zz.index("`") + 2 :]
                )
            elif zz.index("`") + 1 < len(zz) and re_uppercase.search(
                zz[zz.index("`") + 1]
            ):
                zz = (
                    zz[: zz.index("`") + 1]
                    + ""
                    + zz[zz.index("`") + 1]
                    + "&#x0301;"
                    + zz[zz.index("`") + 2 :]
                )
            zz = zz[: zz.index("`")] + zz[zz.index("`") + 1 :]

    return zz


def htmlformat(s, **kwargs):
    res = ""
    for run in parse_4s_elem(s):
        if run[0] == "":
            res += htmlrepl(run[1])
        if run[0] == "em":
            res += "<em>" + htmlrepl(run[1]) + "</em>"
        if run[0] == "img":
            imgfile, w, h = parseimg(
                run[1],
                dimensions="pixels",
                targetdir=kwargs.get("targetdir"),
                tmp_dir=kwargs.get("tmp_dir"),
            )
            if os.path.isfile(imgfile):
                # with open(imgfile, 'rb') as f:
                #     imgdata = f.read()
                # imgfile = 'data:image/{ext};base64,{b64}'.format(
                #     ext=os.path.splitext(imgfile)[-1][1:],
                #     b64=base64.b64encode(imgdata))
                uploaded_image = im.upload_image(imgfile, title=imgfile)
                imgfile = uploaded_image.link

            res += '<img{}{} src="{}"/>'.format(
                "" if w == -1 else " width={}".format(w),
                "" if h == -1 else " height={}".format(h),
                imgfile,
            )
    return res


def htmlyapper(e, **kwargs):
    if isinstance(e, basestring):
        return html_element_layout(e, **kwargs)
    elif isinstance(e, list):
        if not any(isinstance(x, list) for x in e):
            return html_element_layout(e, **kwargs)
        else:
            return "\n".join([html_element_layout(x, **kwargs) for x in e])


def html_element_layout(e, **kwargs):
    res = ""
    if isinstance(e, basestring):
        res = htmlformat(e, **kwargs)
        return res
    if isinstance(e, list):
        res = "\n".join(
            [
                "{}. {}".format(en + 1, html_element_layout(x, **kwargs))
                for en, x in enumerate(e)
            ]
        )
        return res


def md5(s):
    return hashlib.md5(s).hexdigest()


def get_chal(lj, passwd):
    chal = None
    chal = retry_wrapper(lj.getchallenge)["challenge"]
    response = md5(
        chal.encode("utf8") + md5(passwd.encode("utf8")).encode("utf8")
    )
    return (chal, response)


def find_heading(structure):
    h_id = -1
    for e, x in enumerate(structure):
        if x[0] == "ljheading":
            return (e, x)
        elif x[0] == "heading":
            h_id = e
    if h_id >= 0:
        return (h_id, structure[h_id])
    return None


def find_tour(structure):
    for e, x in enumerate(structure):
        if x[0] == "section":
            return (e, x)
    return None


def split_into_tours(structure, general_impression=False):
    result = []
    current = []
    mode = "meta"
    for _, element in enumerate(structure):
        if element[0] != "Question":
            if mode == "meta":
                current.append(element)
            elif element[0] == "section":
                result.append(current)
                current = [element]
                mode = "meta"
            else:
                current.append(element)
        else:
            if mode == "meta":
                current.append(element)
                mode = "questions"
            else:
                current.append(element)
    result.append(current)
    globalheading = find_heading(result[0])[1][1]
    globalsep = "." if not globalheading.endswith(".") else ""
    currentheading = result[0][find_heading(result[0])[0]][1]
    result[0][find_heading(result[0])[0]][1] += "{} {}".format(
        "." if not currentheading.endswith(".") else "",
        find_tour(result[0])[1][1],
    )
    for tour in result[1:]:
        if not find_heading(tour):
            tour.insert(
                0,
                [
                    "ljheading",
                    "{}{} {}".format(
                        globalheading, globalsep, find_tour(tour)[1][1]
                    ),
                ],
            )
    if general_impression:
        result.append(
            [
                [
                    "ljheading",
                    "{}{} Общие впечатления".format(globalheading, globalsep),
                ],
                [
                    "meta",
                    "В комментариях к этому посту можно "
                    "поделиться общими впечатлениями от вопросов.",
                ],
            ]
        )
    return result


def lj_process(structure, args, **kwargs):
    final_structure = [{"header": "", "content": ""}]
    i = 0
    heading = ""
    ljheading = ""

    def yapper(x):
        return htmlyapper(x, **kwargs)

    while i < len(structure) and structure[i][0] != "Question":
        if structure[i][0] == "heading":
            final_structure[0]["content"] += "<center>{}</center>".format(
                yapper(structure[i][1])
            )
            heading = yapper(structure[i][1])
        if structure[i][0] == "ljheading":
            # final_structure[0]['header'] = structure[i][1]
            ljheading = yapper(structure[i][1])
        if structure[i][0] == "date":
            final_structure[0]["content"] += "\n<center>{}</center>".format(
                yapper(structure[i][1])
            )
        if structure[i][0] == "editor":
            final_structure[0]["content"] += "\n<center>{}</center>".format(
                yapper(structure[i][1])
            )
        if structure[i][0] == "meta":
            final_structure[0]["content"] += "\n{}".format(
                yapper(structure[i][1])
            )
        i += 1

    if ljheading != "":
        final_structure[0]["header"] = ljheading
    else:
        final_structure[0]["header"] = heading

    for element in structure[i:]:
        if element[0] == "Question":
            final_structure.append(
                {
                    "header": "Вопрос {}".format(
                        element[1]["number"]
                        if "number" in element[1]
                        else gui_compose.counter
                    ),
                    "content": html_format_question(element[1], **kwargs),
                }
            )
        if element[0] == "meta":
            final_structure.append(
                {"header": "", "content": yapper(element[1])}
            )

    if not final_structure[0]["content"]:
        final_structure[0]["content"] = "Вопросы в комментариях."
    if debug:
        with codecs.open("lj.debug", "w", "utf8") as f:
            f.write(log_wrap(final_structure))
    return final_structure


def _lj_post(lj, stru, args, edit=False, add_params=None):

    now = datetime.datetime.now()
    year = now.strftime("%Y")
    month = now.strftime("%m")
    day = now.strftime("%d")
    hour = now.strftime("%H")
    minute = now.strftime("%M")

    chal, response = get_chal(lj, args.password)

    params = {
        "username": args.login,
        "auth_method": "challenge",
        "auth_challenge": chal,
        "auth_response": response,
        "subject": stru["header"],
        "event": stru["content"],
        "year": year,
        "mon": month,
        "day": day,
        "hour": hour,
        "min": minute,
    }
    if edit:
        params["itemid"] = stru["itemid"]
    if add_params:
        params.update(add_params)

    try:
        post = retry_wrapper(lj.editevent if edit else lj.postevent, [params])
        logger.info("Edited a post" if edit else "Created a post")
        logger.debug(log_wrap(post))
        time.sleep(5)
    except Exception as e:
        sys.stderr.write(
            "Error issued by LJ API: {}".format(traceback.format_exc(e))
        )
        sys.exit(1)
    return post


def _lj_comment(lj, stru, args):
    chal, response = get_chal(lj, args.password)
    params = {
        "username": args.login,
        "auth_method": "challenge",
        "auth_challenge": chal,
        "auth_response": response,
        "journal": stru["journal"],
        "ditemid": stru["ditemid"],
        "parenttalkid": 0,
        "body": stru["content"],
        "subject": stru["header"],
    }
    try:
        comment = retry_wrapper(lj.addcomment, [params])
    except Exception as e:
        sys.stderr.write(
            "Error issued by LJ API: {}".format(traceback.format_exc(e))
        )
        sys.exit(1)
    logger.info("Added a comment")
    logger.debug(log_wrap(comment))
    time.sleep(random.randint(5, 7))


def lj_post(stru, args, edit=False):

    lj = ServerProxy("http://www.livejournal.com/interface/xmlrpc").LJ.XMLRPC

    add_params = {}
    community = args.community
    if community:
        add_params["usejournal"] = community
    elif args.security:
        add_params["security"] = "usemask"
        add_params["allowmask"] = (
            "1" if args.security == "friends" else args.security
        )
    else:
        add_params["security"] = "private"

    journal = community if community else args.login

    post = _lj_post(lj, stru[0], args, edit=edit, add_params=add_params)

    comments = stru[1:]

    if not comments:
        return post

    for comment in stru[1:]:
        comment["ditemid"] = post["ditemid"]
        comment["journal"] = journal
        _lj_comment(lj, comment, args)

    return post


def baseyapper(e, **kwargs):
    if isinstance(e, basestring):
        return base_element_layout(e, **kwargs)
    elif isinstance(e, list):
        if not any(isinstance(x, list) for x in e):
            return base_element_layout(e, **kwargs)
        else:
            return "\n".join([base_element_layout(x, **kwargs) for x in e])


def baseformat(s, **kwargs):
    res = ""
    for run in parse_4s_elem(s):
        if run[0] == "":
            res += run[1].replace("\n", "\n   ")
        if run[0] == "em":
            res += run[1]
        if run[0] == "img":
            imgfile, w, h = parseimg(
                run[1],
                dimensions="pixels",
                targetdir=kwargs.get("targetdir"),
                tmp_dir=kwargs.get("tmp_dir"),
            )
            if os.path.isfile(imgfile):
                pil_image = Image.open(imgfile)
                w_orig, h_orig = pil_image.size
                if w_orig != w or h_orig != h:
                    print("resizing image {}".format(imgfile))
                    pil_image = pil_image.resize(
                        (int(w), int(h)), resample=Image.LANCZOS
                    )
                    bn, ext = os.path.splitext(imgfile)
                    resized_fn = "{}_resized.png".format(bn)
                    pil_image.save(resized_fn, "PNG")
                    to_upload = resized_fn
                else:
                    to_upload = imgfile
                print("uploading {}...".format(to_upload))
                uploaded_image = im.upload_image(to_upload, title=to_upload)
                imglink = uploaded_image.link
                print("the link for {} is {}...".format(to_upload, imglink))
            else:
                imglink = imgfile
            res += "(pic: {})".format(imglink)
    while res.endswith("\n"):
        res = res[:-1]
    return res


def base_element_layout(e, **kwargs):
    res = ""
    if isinstance(e, basestring):
        res = baseformat(e, **kwargs)
        return res
    if isinstance(e, list):
        res = "\n".join(
            [
                "   {}. {}".format(i + 1, base_element_layout(x, **kwargs))
                for i, x in enumerate(e)
            ]
        )
    return res


BASE_MAPPING = {
    "section": "Тур",
    "heading": "Чемпионат",
    "editor": "Редактор",
    "meta": "Инфо",
}
re_date_sep = re.compile(" [—–-] ")


def wrap_date(s):
    s = s.strip()
    parsed = dateparser.parse(s)
    if isinstance(parsed, datetime.datetime):
        parsed = parsed.date()
    if not parsed:
        print(
            "unable to parse date {}, setting to default 2010-01-01".format(s)
        )
        return datetime.date(2010, 1, 1).strftime("%d-%b-%Y")
    if parsed > datetime.date.today():
        parsed = parsed.replace(year=parsed.year - 1)
    formatted = parsed.strftime("%d-%b-%Y")
    return formatted


def base_format_element(pair, **kwargs):
    if pair[0] == "Question":
        return base_format_question(pair[1], **kwargs)
    if pair[0] in BASE_MAPPING:
        return "{}:\n{}\n\n".format(
            BASE_MAPPING[pair[0]], baseyapper(pair[1], **kwargs)
        )
    elif pair[0] == "date":
        re_search = re_date_sep.search(pair[1])
        if re_search:
            gr0 = re_search.group(0)
            dates = pair[1].split(gr0)
            return "Дата:\n{} - {}\n\n".format(
                wrap_date(dates[0]), wrap_date(dates[-1])
            )
        else:
            return "Дата:\n{}\n\n".format(wrap_date(pair[1]))


def check_if_zero(Question):
    number = Question.get("number")
    if number is None:
        return False
    if isinstance(number, int) and number == 0:
        return True
    if isinstance(number, str) and number.startswith(("0", "Размин")):
        return True
    return False


def output_base(structure, outfile, args, **kwargs):
    def _baseyapper(x):
        return baseyapper(x, **kwargs)

    def _get_last_value(dct, key):
        if isinstance(dct[key], list):
            return dct[key][-1]
        return dct[key]

    def _add_to_dct(dct, key, to_add):
        if isinstance(dct[key], list):
            dct[key][-1] += to_add
        else:
            dct[key] += to_add

    result = []
    lasttour = 0
    zeroq = 1
    for i, pair in enumerate(structure):
        if pair[0] == "section":
            lasttour = i
        while (
            pair[0] == "meta"
            and (i + 1) < len(structure)
            and structure[i + 1][0] == "meta"
        ):
            pair[1] += "\n{}".format(structure[i + 1][1])
            structure.pop(i + 1)
        if pair[0] == "Question" and check_if_zero(pair[1]):
            tourheader = "Нулевой вопрос {}".format(zeroq)
            zeroq += 1
            pair[1]["number"] = 1
            structure.insert(lasttour, structure.pop(i))
            structure.insert(lasttour, ["section", tourheader])
    for pair in structure:
        if pair[0] == "Question" and "nezachet" in pair[1]:
            field = "zachet" if "zachet" in pair[1] else "answer"
            last_val = _get_last_value(pair[1], field)
            nezachet = _baseyapper(pair[1].pop("nezachet"))
            to_add = "{}\n   Незачёт: {}".format(
                "." if not last_val.endswith(".") else "", nezachet
            )
            _add_to_dct(pair[1], field, to_add)
        if pair[0] == "editor":
            pair[1] = re.sub(re_editors, "", pair[1])
            logger.info('Поле "Редактор" было автоматически изменено.')
        res = base_format_element(pair, **kwargs)
        if res:
            result.append(res)
    text = "".join(result)
    with codecs.open(outfile, "w", "utf8") as f:
        f.write(text)
    logger.info("Output: {}".format(outfile))
    if args.clipboard:
        pyperclip.copy(text)


def base_format_question(q, **kwargs):
    def _baseyapper(x):
        return baseyapper(x, **kwargs)

    if "setcounter" in q:
        gui_compose.counter = int(q["setcounter"])
    res = "Вопрос {}:\n{}\n\n".format(
        gui_compose.counter if "number" not in q else q["number"],
        _baseyapper(q["question"]),
    )
    if "number" not in q:
        gui_compose.counter += 1
    res += "Ответ:\n{}\n\n".format(_baseyapper(q["answer"]))
    if "zachet" in q:
        res += "Зачет:\n{}\n\n".format(_baseyapper(q["zachet"]))
    if "nezachet" in q:
        res += "Незачет:\n{}\n\n".format(_baseyapper(q["zachet"]))
    if "comment" in q:
        res += "Комментарий:\n{}\n\n".format(_baseyapper(q["comment"]))
    if "source" in q:
        res += "Источник:\n{}\n\n".format(_baseyapper(q["source"]))
    if "author" in q:
        res += "Автор:\n{}\n\n".format(_baseyapper(q["author"]))
    return res


def reddityapper(e, **kwargs):
    if isinstance(e, basestring):
        return reddit_element_layout(e, **kwargs)
    elif isinstance(e, list):
        if not any(isinstance(x, list) for x in e):
            return reddit_element_layout(e, **kwargs)
        else:
            return "  \n".join([reddit_element_layout(x, **kwargs) for x in e])


def redditformat(s, **kwargs):
    res = ""
    for run in parse_4s_elem(s):
        if run[0] == "":
            res += run[1]
        if run[0] == "em":
            res += "_{}_".format(run[1])
        if run[0] == "img":
            imgfile, w, h = parseimg(
                run[1],
                dimensions="ems",
                targetdir=kwargs.get("targetdir"),
                tmp_dir=kwargs.get("tmp_dir"),
            )
            if os.path.isfile(imgfile):
                uploaded_image = im.upload_image(imgfile, title=imgfile)
                imgfile = uploaded_image.link
            else:
                raise Exception("Image not found: {}".format(imgfile))
            res += "[картинка]({})".format(imgfile)
    while res.endswith("\n"):
        res = res[:-1]
    res = res.replace("\n", "  \n")
    return res


def reddit_element_layout(e, **kwargs):
    res = ""
    if isinstance(e, basestring):
        res = redditformat(e, **kwargs)
        return res
    if isinstance(e, list):
        res = "  \n".join(
            [
                "{}\\. {}".format(i + 1, reddit_element_layout(x, **kwargs))
                for i, x in enumerate(e)
            ]
        )
    return res


def reddit_format_element(pair, **kwargs):
    if pair[0] == "Question":
        return reddit_format_question(pair[1], **kwargs)


def reddit_format_question(q, **kwargs):
    if "setcounter" in q:
        gui_compose.counter = int(q["setcounter"])
    res = "__Вопрос {}__: {}  \n".format(
        gui_compose.counter if "number" not in q else q["number"],
        reddityapper(q["question"], **kwargs),
    )
    if "number" not in q:
        gui_compose.counter += 1
    res += "__Ответ:__ >!{}  \n".format(reddityapper(q["answer"], **kwargs))
    if "zachet" in q:
        res += "__Зачёт:__ {}  \n".format(reddityapper(q["zachet"], **kwargs))
    if "nezachet" in q:
        res += "__Незачёт:__ {}  \n".format(
            reddityapper(q["nezachet"], **kwargs)
        )
    if "comment" in q:
        res += "__Комментарий:__ {}  \n".format(
            reddityapper(q["comment"], **kwargs)
        )
    if "source" in q:
        res += "__Источник:__ {}  \n".format(
            reddityapper(q["source"], **kwargs)
        )
    if "author" in q:
        res += "!<\n__Автор:__ {}  \n".format(
            reddityapper(q["author"], **kwargs)
        )
    else:
        res += "!<\n"
    return res


def output_reddit(structure, outfile, args, **kwargs):
    result = []
    for pair in structure:
        res = reddit_format_element(pair, **kwargs)
        if res:
            result.append(res)
    text = "\n\n".join(result)
    with codecs.open(outfile, "w", "utf8") as f:
        f.write(text)
    logger.info("Output: {}".format(outfile))


def tex_format_question(q, **kwargs):
    yapper = texyapper
    if "setcounter" in q:
        gui_compose.counter = int(q["setcounter"])
    res = (
        "\n\n\\begin{{minipage}}{{\\textwidth}}\\raggedright\n"
        "\\textbf{{Вопрос {}.}} {} \\newline".format(
            gui_compose.counter if "number" not in q else q["number"],
            yapper(q["question"], **kwargs),
        )
    )
    if "number" not in q:
        gui_compose.counter += 1
    res += "\n\\textbf{{Ответ: }}{} \\newline".format(
        yapper(q["answer"], **kwargs)
    )
    if "zachet" in q:
        res += "\n\\textbf{{Зачёт: }}{} \\newline".format(
            yapper(q["zachet"], **kwargs)
        )
    if "nezachet" in q:
        res += "\n\\textbf{{Незачёт: }}{} \\newline".format(
            yapper(q["nezachet"], **kwargs)
        )
    if "comment" in q:
        res += "\n\\textbf{{Комментарий: }}{} \\newline".format(
            yapper(q["comment"], **kwargs)
        )
    if "source" in q:
        res += "\n\\textbf{{Источник{}: }}{} \\newline".format(
            "и" if isinstance(q["source"], list) else "",
            yapper(q["source"], **kwargs),
        )
    if "author" in q:
        res += "\n\\textbf{{Автор: }}{} \\newline".format(
            yapper(q["author"], **kwargs)
        )
    res += "\n\\end{minipage}\n"
    return res


def texrepl(zz):
    zz = re.sub(r"{", r"\{", zz)
    zz = re.sub(r"}", r"\}", zz)
    zz = re.sub(r"\\(?![\}\{])", r"{\\textbackslash}", zz)
    zz = re.sub("%", "\%", zz)
    zz = re.sub(r"\$", "\$", zz)
    zz = re.sub("#", "\#", zz)
    zz = re.sub("&", "\&", zz)
    zz = re.sub("_", r"\_", zz)
    zz = re.sub(r"\^", r"{\\textasciicircum}", zz)
    zz = re.sub(r"\~", r"{\\textasciitilde}", zz)
    zz = re.sub(r'((\"(?=[ \.\,;\:\?!\)\]]))|("(?=\Z)))', "»", zz)
    zz = re.sub(r'(((?<=[ \.\,;\:\?!\(\[)])")|((?<=\A)"))', "«", zz)
    zz = re.sub('"', "''", zz)

    for match in sorted(
        [x for x in re_scaps.finditer(zz)],
        key=lambda x: len(x.group(2)),
        reverse=True,
    ):
        zz = zz.replace(
            match.group(2), "\\tsc{" + match.group(2).lower() + "}"
        )

    # while re_scaps.search(zz):
    #     zz = zz.replace(re_scaps.search(zz).group(2),
    #         '\\tsc{'+re_scaps.search(zz).group(2).lower()+'}')

    torepl = [x.group(0) for x in re.finditer(re_url, zz)]
    for s in range(len(torepl)):
        item = torepl[s]
        while item[-1] in typotools.PUNCTUATION:
            item = item[:-1]
        while (
            item[-1] in typotools.CLOSING_BRACKETS
            and typotools.find_matching_opening_bracket(item, -1) is None
        ):
            item = item[:-1]
        while item[-1] in typotools.PUNCTUATION:
            item = item[:-1]
        torepl[s] = item
    torepl = sorted(set(torepl), key=len, reverse=True)
    hashurls = {}
    for s in torepl:
        hashurls[s] = hashlib.md5(s.encode("utf8")).hexdigest()
    for s in sorted(hashurls, key=len, reverse=True):
        zz = zz.replace(s, hashurls[s])
    hashurls = {v: k for k, v in hashurls.items()}
    for s in sorted(hashurls):
        zz = zz.replace(
            s, "\\url{{{}}}".format(hashurls[s].replace("\\\\", "\\"))
        )

    # debug_print('URLS FOR REPLACING: ' +
    #             pprint.pformat(torepl).decode('unicode_escape'))
    # while len(torepl)>0:
    #     s = torepl[0]
    #     debug_print('STRING BEFORE REPLACEMENT: {}'.format(zz))
    #     zz = zz.replace(s, '\\url{'+s+'}')
    #     debug_print('STRING AFTER REPLACEMENT: {}'.format(zz))
    #     torepl.pop(0)

    zz = zz.replace(" — ", "{\\Hair}—{\\hair}")

    while "`" in zz:
        if zz.index("`") + 1 >= len(zz):
            zz = zz.replace("`", "")
        else:
            if zz.index("`") + 2 < len(zz) and re.search(
                r"\s", zz[zz.index("`") + 2]
            ):
                zz = zz[: zz.index("`") + 2] + "" + zz[zz.index("`") + 2 :]
            if zz.index("`") + 1 < len(zz) and re_lowercase.search(
                zz[zz.index("`") + 1]
            ):
                zz = (
                    zz[: zz.index("`") + 1]
                    + ""
                    + zz[zz.index("`") + 1]
                    + "\u0301"
                    + zz[zz.index("`") + 2 :]
                )
            elif zz.index("`") + 1 < len(zz) and re_uppercase.search(
                zz[zz.index("`") + 1]
            ):
                zz = (
                    zz[: zz.index("`") + 1]
                    + ""
                    + zz[zz.index("`") + 1]
                    + "\u0301"
                    + zz[zz.index("`") + 2 :]
                )
            zz = zz[: zz.index("`")] + zz[zz.index("`") + 1 :]

    return zz


def texformat(s, **kwargs):
    res = ""
    for run in parse_4s_elem(s):
        if run[0] == "":
            res += texrepl(run[1])
        if run[0] == "em":
            res += "\\emph{" + texrepl(run[1]) + "}"
        if run[0] == "img":
            imgfile, w, h = parseimg(
                run[1],
                dimensions="ems",
                tmp_dir=kwargs.get("tmp_dir"),
                targetdir=kwargs.get("targetdir"),
            )
            res += (
                "\\includegraphics"
                + "[width={}{}]".format(
                    "10em" if w == -1 else "{}em".format(w),
                    ", height={}em".format(h) if h != -1 else "",
                )
                + "{"
                + imgfile
                + "}"
            )
    while res.endswith("\n"):
        res = res[:-1]
    res = res.replace("\n", "  \\newline \n")
    return res


def texyapper(e, **kwargs):
    if isinstance(e, basestring):
        return tex_element_layout(e, **kwargs)
    elif isinstance(e, list):
        if not any(isinstance(x, list) for x in e):
            return tex_element_layout(e, **kwargs)
        else:
            return "  \n".join([tex_element_layout(x, **kwargs) for x in e])


def tex_element_layout(e, **kwargs):
    res = ""
    if isinstance(e, basestring):
        res = texformat(e, **kwargs)
        return res
    if isinstance(e, list):
        res = """
\\begin{{compactenum}}
{}
\\end{{compactenum}}
""".format(
            "\n".join(
                [
                    "\\item {}".format(tex_element_layout(x, **kwargs))
                    for x in e
                ]
            )
        )
    return res


def gui_compose(largs, sourcedir=None):
    global args
    global console_mode
    args = largs
    global debug
    global logger
    global retry_wrapper
    assert sourcedir is not None

    logger = logging.getLogger("composer")
    logger.setLevel(logging.DEBUG)
    fh = logging.FileHandler("composer.log", encoding="utf8")
    fh.setLevel(logging.DEBUG)
    ch = logging.StreamHandler()
    if args.debug:
        ch.setLevel(logging.DEBUG)
    else:
        ch.setLevel(logging.INFO)
    formatter = logging.Formatter("%(asctime)s | %(message)s")
    fh.setFormatter(formatter)
    ch.setFormatter(formatter)
    logger.addHandler(fh)
    logger.addHandler(ch)

    retry_wrapper = retry_wrapper_factory(logger)

    if args.debug:
        debug = True

    argsdict = vars(args)
    logger.debug(log_wrap(argsdict))

    if args.filename and args.filetype:
        if args.filetype == "lj":
            if args.login and args.password:
                console_mode = True
        else:
            console_mode = True

    ld = get_lastdir()
    if args.filename:
        if isinstance(args.filename, list):
            ld = os.path.dirname(os.path.abspath(args.filename[0]))
        else:
            ld = os.path.dirname(os.path.abspath(args.filename))
    set_lastdir(ld)
    if not args.filename:
        print("No file specified.")
        sys.exit(1)

    if isinstance(args.filename, list):
        if not args.merge:
            for fn in args.filename:
                targetdir = os.path.dirname(os.path.abspath(fn))
                filename = os.path.basename(os.path.abspath(fn))
                process_file_wrapper(filename, sourcedir, targetdir)
        else:
            targetdir = os.path.dirname(os.path.abspath(args.filename[0]))
            process_file_wrapper(args.filename, sourcedir, targetdir)
    else:
        targetdir = os.path.dirname(os.path.abspath(args.filename))
        filename = os.path.basename(os.path.abspath(args.filename))
        process_file_wrapper(filename, sourcedir, targetdir)


def process_file_wrapper(filename, sourcedir, targetdir):
    resourcedir = os.path.join(sourcedir, "resources")
    with make_temp_directory(dir=get_chgksuite_dir()) as tmp_dir:
        for fn in [
            args.docx_template,
            os.path.join(resourcedir, "fix-unnumbered-sections.sty"),
            args.tex_header,
        ]:
            shutil.copy(fn, tmp_dir)
        process_file(filename, tmp_dir, sourcedir, targetdir)


def parse_filepath(filepath):
    with codecs.open(filepath, "r", "utf8") as input_file:
        input_text = input_file.read()
    input_text = input_text.replace("\r", "")
    return parse_4s(input_text, randomize=args.randomize)


def make_merged_filename(filelist):
    filelist = [os.path.splitext(os.path.basename(x))[0] for x in filelist]
    prefix = os.path.commonprefix(filelist)
    suffix = "_".join(x[len(prefix) :] for x in filelist)
    return prefix + suffix


def generate_navigation(strus):
    titles = [x[0][0]["header"].split(". ")[-1] for x in strus]
    urls = [x[1]["url"] for x in strus]
    result = []
    for i in range(len(titles)):
        inner = []
        for j in range(len(urls)):
            inner.append(
                titles[j]
                if j == i
                else '<a href="{}">{}</a>'.format(urls[j], titles[j])
            )
        result.append(" | ".join(inner))
    return result


def find_min_context_index(structure):
    types_ = [x[0] for x in structure]
    try:
        min_section = types_.index("section")
    except ValueError:
        min_section = None
    try:
        min_question = types_.index("Question")
    except ValueError:
        min_question = None
    if min_section is not None and min_question is not None:
        return min(min_section, min_question)
    elif min_section is not None:
        return min_section
    else:
        return min_question


def backtick_replace(el):
    while "`" in el:
        if el.index("`") + 1 >= len(el):
            el = el.replace("`", "")
        else:
            if el.index("`") + 2 < len(el) and re.search(
                r"\s", el[el.index("`") + 2]
            ):
                el = el[: el.index("`") + 2] + "" + el[el.index("`") + 2 :]
            if el.index("`") + 1 < len(el) and re_lowercase.search(
                el[el.index("`") + 1]
            ):
                el = (
                    el[: el.index("`") + 1]
                    + ""
                    + el[el.index("`") + 1]
                    + "\u0301"
                    + el[el.index("`") + 2 :]
                )
            elif el.index("`") + 1 < len(el) and re_uppercase.search(
                el[el.index("`") + 1]
            ):
                el = (
                    el[: el.index("`") + 1]
                    + ""
                    + el[el.index("`") + 1]
                    + "\u0301"
                    + el[el.index("`") + 2 :]
                )
            el = el[: el.index("`")] + el[el.index("`") + 1 :]
    return el


class DocxExporter(object):
    def __init__(self, structure, args, dir_kwargs):
        self.structure = structure
        self.args = args
        self.dir_kwargs = dir_kwargs
        self.qcount = 0

    def _docx_format(self, *args):
        return self.docx_format(*args, **self.dir_kwargs)

    def docx_format(self, el, para, whiten, **kwargs):
        if isinstance(el, list):

            if len(el) > 1 and isinstance(el[1], list):
                self.docx_format(el[0], para, whiten, **kwargs)
                licount = 0
                for li in el[1]:
                    licount += 1

                    para.add_run("\n{}. ".format(licount))
                    self.docx_format(li, para, whiten, **kwargs)
            else:
                licount = 0
                for li in el:
                    licount += 1

                    para.add_run("\n{}. ".format(licount))
                    self.docx_format(li, para, whiten, **kwargs)

        if isinstance(el, basestring):
            logger.debug("parsing element {}:".format(log_wrap(el)))

            el = backtick_replace(el)
            parsed = parse_4s_elem(el)
            images_exist = False

            for run in parsed:
                if run[0] == "img":
                    images_exist = True

            for run in parse_4s_elem(el):
                if run[0] == "":
                    r = para.add_run(run[1])
                    if whiten and not args.nospoilers:
                        r.style = "Whitened"

                elif run[0] == "em":
                    r = para.add_run(run[1])
                    r.italic = True
                    if whiten and not args.nospoilers:
                        r.style = "Whitened"

                elif run[0] == "sc":
                    r = para.add_run(run[1])
                    r.small_caps = True
                    if whiten and not args.nospoilers:
                        r.style = "Whitened"

                elif run[0] == "img":
                    imgfile, width, height = parseimg(
                        run[1],
                        dimensions="inches",
                        tmp_dir=kwargs.get("tmp_dir"),
                        targetdir=kwargs.get("targetdir"),
                    )
                    r = para.add_run("\n")
                    r.add_picture(
                        imgfile, width=Inches(width), height=Inches(height)
                    )
                    r.add_text("\n")

    def add_question(self, element):
        q = element[1]
        p = self.doc.add_paragraph()
        p.paragraph_format.keep_together = True
        if "number" not in q:
            self.qcount += 1
        if "setcounter" in q:
            self.qcount = int(q["setcounter"])
        p.add_run(
            "Вопрос {}. ".format(
                qcount if "number" not in q else q["number"]
            )
        ).bold = True

        if "handout" in q:
            p.add_run("\n[Раздаточный материал: ")
            self._docx_format(q["handout"], p, WHITEN["handout"])
            p.add_run("\n]")
        if not args.noparagraph:
            p.add_run("\n")

        self._docx_format(q["question"], p, False)
        p = self.doc.add_paragraph()
        p.paragraph_format.keep_together = True

        if not args.noanswers:
            if not args.no_line_break:
                p = self.doc.add_paragraph()
                p.paragraph_format.keep_together = True
            p.add_run("Ответ: ").bold = True
            self._docx_format(q["answer"], p, True)

            for field in [
                "zachet",
                "nezachet",
                "comment",
                "source",
                "author",
            ]:
                if field in q:
                    if field == "source":
                        p = self.doc.add_paragraph()
                        p.paragraph_format.keep_together = True
                    else:
                        p.add_run("\n")
                    if field == "source" and isinstance(
                        q[field], list
                    ):
                        p.add_run("Источники: ").bold = True
                    else:
                        p.add_run(FIELDS[field]).bold = True
                    self._docx_format(q[field], p, WHITEN[field])

        self.doc.add_paragraph()
        if not args.one_line_break:
            self.doc.add_paragraph()

    def export(self, outfilename):
        logger.debug(args.docx_template)
        self.doc = Document(args.docx_template)
        logger.debug(log_wrap(self.structure))

        for element in self.structure:
            if element[0] == "meta":
                p = self.doc.add_paragraph()
                self._docx_format(element[1], p, False)
                self.doc.add_paragraph()

            if element[0] in ["editor", "date", "heading", "section"]:
                para = self.doc.add_paragraph(element[1])
                para.alignment = 1
                para.paragraph_format.keep_with_next = True
                para.add_run("\n")

            if element[0] == "Question":
                self.add_question(element)

        self.doc.save(outfilename)
        logger.info("Output: {}".format(outfilename))


class PptxExporter(object):
    def __init__(self, structure, args, dir_kwargs):
        self.structure = structure
        self.args = args
        self.config_path = os.path.abspath(args.pptx_config)
        with open(self.config_path) as f:
            self.c = toml.load(f)
        self.dir_kwargs = dir_kwargs
        self.qcount = 0

    @staticmethod
    def add_editor_info(tb, editor, meta, slide):
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = editor
        for elem in meta:
            p = tf.add_paragraph()
            self.pptx_format(self.pptx_process_text(elem[1]), para, tf, slide)

    def get_textbox_qnumber(self, slide):
        return self.get_textbox(
            slide,
            left=PptxInches(self.c["number_textbox"]["left"]),
            width=PptxInches(self.c["number_textbox"]["width"]),
            height=PptxInches(self.c["number_textbox"]["height"]),
        )

    def get_textbox(self, slide, left=None, top=None, width=None, height=None):
        if left is None:
            left = PptxInches(self.c["textbox"]["left"])
        if top is None:
            top = PptxInches(self.c["textbox"]["top"])
        if width is None:
            width = PptxInches(self.c["textbox"]["width"])
        if height is None:
            height = PptxInches(self.c["textbox"]["height"])
        textbox = slide.shapes.add_textbox(left, top, width, height)
        return textbox

    def pptx_format(self, el, para, tf, slide):
        if isinstance(el, list):
            if len(el) > 1 and isinstance(el[1], list):
                self.pptx_format(el[0], para, tf, slide)
                licount = 0
                for li in el[1]:
                    licount += 1
                    p = self.init_paragraph(tf)
                    r = p.add_run()
                    r.text = "{}. ".format(licount)
                    self.pptx_format(li, p, tf, slide)
            else:
                licount = 0
                for li in el:
                    licount += 1
                    p = self.init_paragraph(tf)
                    r = p.add_run()
                    r.text = "{}. ".format(licount)
                    self.pptx_format(li, p, tf, slide)

        if isinstance(el, basestring):
            logger.debug("parsing element {}:".format(log_wrap(el)))
            el = backtick_replace(el)

            for run in parse_4s_elem(el):
                if run[0] in ("", "sc"):
                    r = para.add_run()
                    r.text = run[1]

                elif run[0] == "em":
                    r = para.add_run()
                    r.text = run[1]
                    r.italic = True

                elif run[0] == "img":
                    imgfile, width, height = parseimg(
                        run[1],
                        dimensions="inches",
                        tmp_dir=self.dir_kwargs.get("tmp_dir"),
                        targetdir=self.dir_kwargs.get("targetdir"),
                    )
                    slide.shapes.add_picture(
                        imgfile,
                        left=PptxInches(self.c["picture"]["left"]),
                        top=PptxInches(self.c["picture"]["top"]),
                        width=PptxInches(width),
                        height=PptxInches(height),
                    )
                    para = self.init_paragraph(tf)

    def pptx_process_text(self, s):
        if isinstance(s, list):
            for i in range(len(s)):
                s[i] = self.pptx_process_text(s[i])
            return s
        s = s.replace("\u0301", "")
        if "Раздат" not in s:
            while "[" in s and "]" in s:
                s = re.sub("\[.+?\]", "", s)
        s = re.sub(" +", " ", s)
        for punct in (".", ",", "!", "?", ":"):
            s = s.replace(" " + punct, punct)
        return s

    def process_header(self, header):
        if "editor" in header:
            index_ = [x[0] for x in header].index("editor")
            editor_info = header[index_:]
            header = header[:index_]
        else:
            editor_info = []
        heading = [x for x in header if x[0] == "heading"]
        ljheading = [x for x in header if x[0] == "ljheading"]
        title_text = heading or ljheading
        date_text = [x for x in header if x[0] == "date"]
        if title_text:
            slide = self.prs.slides.add_slide(self.TITLE_SLIDE)
            title = slide.shapes.title
            title.text = title_text[0][1]
            if date_text:
                subtitle = slide.placeholders[1]
                subtitle.text = date_text[0][1]
        if editor_info:
            slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
            textbox = self.get_textbox(slide)
            editor = [x[1] for x in editor_info if x[0] == "editor"][0]
            meta = [x for x in editor_info if x[0] == "meta"]
            self.add_editor_info(textbox, editor, meta, slide)

    def set_question_number(self, slide, number):
        qntextbox = self.get_textbox_qnumber(slide)
        qtf = qntextbox.text_frame
        qtf_p = self.init_paragraph(qtf)
        qtf_r = qtf_p.add_run()
        qtf_r.text = number

    def process_question(self, q):
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        textbox = self.get_textbox(slide)
        tf = textbox.text_frame
        tf.word_wrap = True
        if "number" not in q:
            self.qcount += 1
        if "setcounter" in q:
            self.qcount = int(q["setcounter"])
        number = str(self.qcount if "number" not in q else q["number"])
        self.set_question_number(slide, number)
        # p = self.init_paragraph(tf)
        # p.alignment = PP_ALIGN.CENTER
        # r = p.add_run()
        # r.text = "Вопрос {}. ".format(
        #     self.qcount if "number" not in q else q["number"]
        # )
        # r.font.bold = True
        p = self.init_paragraph(tf)

        if "handout" in q:
            p = self.init_paragraph(tf)
            p.add_run("[Раздаточный материал: ")
            self.pptx_format(q["handout"], p, tf, slide)
            p = self.init_paragraph(tf)
            p.add_run("]")

        question_text = self.pptx_process_text(q["question"])
        p.font.size = Pt(self.determine_size(question_text))
        self.pptx_format(question_text, p, tf, slide)

        if self.c["add_plug"]:
            slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
            self.set_question_number(slide, number)
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        self.set_question_number(slide, number)
        textbox = self.get_textbox(slide)
        tf = textbox.text_frame
        tf.word_wrap = True

        p = self.init_paragraph(tf)
        r = p.add_run()
        r.text = "Ответ: "
        r.font.bold = True
        self.pptx_format(self.pptx_process_text(q["answer"]), p, tf, slide)
        if self.c["add_comment"]:
            comment_text = self.pptx_process_text(q["comment"])
            p = self.init_paragraph(tf, text=comment_text)
            r = p.add_run()
            r.text = "Комментарий: "
            r.font.bold = True
            self.pptx_format(comment_text, p, tf, slide)

    def determine_size(self, text):
        for element in self.c["text_size_grid"]["elements"]:
            if len(text) <= element["length"]:
                return element["size"]
        return self.c["text_size_grid"]["smallest"]

    def init_paragraph(self, text_frame, text=None):
        p = text_frame.add_paragraph()
        p.font.name = self.c["font"]["name"]
        size = self.c["text_size_grid"]["default"]
        if text:
            size = self.determine_size(text)
        p.font.size = Pt(size)
        return p

    def export(self, outfilename):
        self.outfilename = outfilename
        wd = os.getcwd()
        os.chdir(os.path.dirname(self.config_path))
        template = os.path.abspath(self.c["template_path"])
        os.chdir(wd)
        self.prs = Presentation(template)
        self.TITLE_SLIDE = self.prs.slide_layouts[0]
        self.BLANK_SLIDE = self.prs.slide_layouts[6]
        min_content_index = find_min_context_index(self.structure)
        header = self.structure[:min_content_index]
        content = self.structure[min_content_index:]
        self.process_header(header)
        slide = None
        for element in content:
            if element[0] == "section":
                slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
                textbox = self.get_textbox(slide)
                tf = textbox.text_frame
                tf.word_wrap = True
                p = self.init_paragraph(tf)
                p.text = element[1]
                p.font.size = Pt(self.c["text_size_grid"]["section"])
            if element[0] == "editor":
                if slide is None:
                    slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
                    textbox = self.get_textbox(slide)
                    tf = textbox.text_frame
                    tf.word_wrap = True
                p = self.init_paragraph(tf)
                p.text = element[1]
            if element[0] == "meta":
                p = self.init_paragraph(tf, text=element[1])
                self.pptx_format(element[1], p, tf, slide)
            if element[0] == "Question":
                self.process_question(element[1])
        self.prs.save(outfilename)


def process_file(filename, tmp_dir, sourcedir, targetdir):
    global args
    dir_kwargs = dict(tmp_dir=tmp_dir, targetdir=targetdir)

    if isinstance(filename, list):
        structure = []
        for x in filename:
            structure.extend(parse_filepath(os.path.join(targetdir, x)))
        filename = make_merged_filename(filename)
    else:
        structure = parse_filepath(os.path.join(targetdir, filename))

    if args.debug:
        debug_fn = os.path.join(
            targetdir,
            make_filename(os.path.basename(filename), "dbg", nots=args.nots),
        )
        with codecs.open(debug_fn, "w", "utf8") as output_file:
            output_file.write(
                json.dumps(structure, indent=2, ensure_ascii=False)
            )

    if not args.filetype:
        print("Filetype not specified.")
        sys.exit(1)
    logger.info(
        "Exporting to {}, spoilers are {}...\n".format(
            args.filetype, "off" if args.nospoilers else "on"
        )
    )

    if args.filetype == "docx":

        outfilename = os.path.join(
            targetdir, make_filename(filename, "docx", nots=args.nots)
        )
        exporter = DocxExporter(structure, args, dir_kwargs)
        exporter.export(outfilename)

    if args.filetype == "tex":

        outfilename = os.path.join(
            tmp_dir, make_filename(filename, "tex", nots=args.nots)
        )

        gui_compose.counter = 1

        title = ""
        date = ""
        gui_compose.tex = """\\input{@header}
\\begin{document}
""".replace(
            "@header", os.path.basename(args.tex_header)
        )
        firsttour = True
        for element in structure:
            if element[0] == "heading":
                gui_compose.tex += (
                    "\n{{\\huge {}}}\n"
                    "\\vspace{{0.8em}}\n".format(
                        tex_element_layout(element[1])
                    )
                )
            if element[0] == "date":
                gui_compose.tex += (
                    "\n{{\\large {}}}\n"
                    "\\vspace{{0.8em}}\n".format(
                        tex_element_layout(element[1])
                    )
                )
            if element[0] in {"meta", "editor"}:
                gui_compose.tex += "\n{}\n\\vspace{{0.8em}}\n".format(
                    tex_element_layout(element[1])
                )
            elif element[0] == "section":
                gui_compose.tex += "\n{}\\section*{{{}}}\n\n".format(
                    "\\clearpage" if not firsttour else "",
                    tex_element_layout(element[1]),
                )
                firsttour = False
            elif element[0] == "Question":
                gui_compose.tex += tex_format_question(
                    element[1], **dir_kwargs
                )

        gui_compose.tex += "\\end{document}"

        with codecs.open(outfilename, "w", "utf8") as outfile:
            outfile.write(gui_compose.tex)
        cwd = os.getcwd()
        os.chdir(tmp_dir)
        subprocess.call(
            shlex.split(
                'xelatex -synctex=1 -interaction=nonstopmode "{}"'.format(
                    outfilename
                )
            )
        )
        os.chdir(cwd)
        pdf_filename = (
            os.path.splitext(os.path.basename(outfilename))[0] + ".pdf"
        )
        logger.info("Output: {}".format(os.path.join(targetdir, pdf_filename)))
        shutil.copy(os.path.join(tmp_dir, pdf_filename), targetdir)
        if args.rawtex:
            shutil.copy(outfilename, targetdir)
            shutil.copy(args.tex_header, targetdir)
            shutil.copy(
                os.path.join(tmp_dir, "fix-unnumbered-sections.sty"), targetdir
            )

    if args.filetype == "lj":

        if not args.community:
            args.community = ""
        if not args.login:
            print("Login not specified.")
            sys.exit(1)
        elif not args.password:
            import getpass

            args.password = getpass.getpass()

        gui_compose.counter = 1
        if args.splittours:
            tours = split_into_tours(structure, general_impression=args.genimp)
            strus = []
            for tour in tours:
                stru = lj_process(tour, args, **dir_kwargs)
                post = lj_post(stru, args)
                strus.append((stru, post))
            if args.navigation:
                navigation = generate_navigation(strus)
                for i, (stru, post) in enumerate(strus):
                    newstru = {
                        "header": stru[0]["header"],
                        "content": stru[0]["content"] + "\n\n" + navigation[i],
                        "itemid": post["itemid"],
                    }
                    lj_post([newstru], args, edit=True)
        else:
            stru = lj_process(structure, args, **dir_kwargs)
            post = lj_post(stru, args)

    if args.filetype == "base":
        gui_compose.counter = 1
        outfilename = os.path.join(
            targetdir, make_filename(filename, "txt", nots=args.nots)
        )
        output_base(structure, outfilename, args, **dir_kwargs)

    if args.filetype == "redditmd":
        gui_compose.counter = 1
        outfilename = os.path.join(
            targetdir, make_filename(filename, "md", nots=args.nots)
        )
        output_reddit(structure, outfilename, args, **dir_kwargs)

    if args.filetype == "pptx":
        outfilename = os.path.join(
            targetdir, make_filename(filename, "pptx", nots=args.nots)
        )
        exporter = PptxExporter(structure, args, dir_kwargs)
        exporter.export(outfilename)

    if not console_mode:
        input("Press Enter to continue...")


def main():
    print("This program was not designed to run standalone.")
    input("Press Enter to continue...")


if __name__ == "__main__":
    main()
